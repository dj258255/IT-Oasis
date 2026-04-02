#!/usr/bin/env python3
"""WikiEngine Portfolio - 15p Vertical A4, dense text + images, 5-step framework"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

IMG = Path("/Users/beomsu/Desktop/gitblog/public/uploads/project/WikiEngine")
OUT = Path("/Users/beomsu/Desktop/gitblog/WikiEngine_Portfolio.pptx")

ACCENT=RGBColor(0x4A,0x8D,0xB8); WHITE=RGBColor(0x1E,0x29,0x3B); BG=RGBColor(0xFF,0xFF,0xFF)
GRAY=RGBColor(0x64,0x74,0x8B); LG=RGBColor(0x94,0xA3,0xB8); RED=RGBColor(0xE7,0x4C,0x3C)
GREEN=RGBColor(0x27,0xAE,0x60); ORANGE=RGBColor(0xE6,0x7E,0x22)

prs = Presentation()
prs.slide_width = Cm(21); prs.slide_height = Cm(29.7)

def sl():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    f = s.background.fill; f.solid(); f.fore_color.rgb = BG
    return s

def tx(s,l,t,w,h,text,sz=10,c=WHITE,b=False,al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Cm(l),Cm(t),Cm(w),Cm(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text=text; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=b; p.alignment=al
    return tf

def pa(tf,text,sz=9,c=WHITE,b=False,sp=2):
    p = tf.add_paragraph(); p.text=text; p.font.size=Pt(sz); p.font.color.rgb=c; p.font.bold=b; p.space_before=Pt(sp)

def sec(s,y,title):
    shape = s.shapes.add_shape(1,Cm(1),Cm(y),Cm(19),Cm(0.65))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
    tx(s,1.2,y+0.03,18,0.55,title,sz=11,c=RGBColor(0xFF,0xFF,0xFF),b=True)
    return y+0.8

def lb(s,y,text,c=ACCENT):
    tx(s,1,y,19,0.4,text,sz=9,c=c,b=True)
    return y+0.45

def im(s,path,l,t,w=None,h=None):
    p=Path(path)
    if not p.exists() or p.suffix.lower() in ('.svg','.gif'): return
    kw={}
    if w: kw['width']=Cm(w)
    if h: kw['height']=Cm(h)
    try: s.shapes.add_picture(str(p),Cm(l),Cm(t),**kw)
    except: pass

M=1  # left margin

# ============ P1: 표지 ============
s=sl()
tx(s,M,4,19,2,"WikiEngine",sz=40,b=True,al=PP_ALIGN.CENTER)
tx(s,M,6.5,19,1,"1,215만 건 검색 엔진의 설계부터 RAG까지",sz=20,c=ACCENT,al=PP_ALIGN.CENTER)
tf=tx(s,2,9,17,6,"",sz=10,c=GRAY,al=PP_ALIGN.CENTER)
pa(tf,"나무위키 + 한국어/영어 위키백과 + 뉴스 + 웹텍스트 + C4 코퍼스",sz=10,c=GRAY,sp=4)
pa(tf,"6개 소스, 12,156,589건, 30개 카테고리, 216만 고유 태그",sz=9,c=LG,sp=4)
pa(tf,"",sz=6)
pa(tf,"Spring Boot · Lucene 10.3 + Nori · MySQL 8 · Redis 3노드",sz=10,c=LG)
pa(tf,"Kafka + Debezium CDC · XGBoost LambdaMART · Spring AI + Gemini",sz=10,c=LG,sp=2)
pa(tf,"Nginx L7 · k6 · Grafana + Prometheus · Docker · Ansible · Next.js",sz=10,c=LG,sp=2)
pa(tf,"",sz=8)
pa(tf,"OCI Free Tier ARM Ampere A1 x 2대 (각 2코어 / 12GB RAM)",sz=9,c=LG)
pa(tf,"2개월 개발 · 26편 기술 블로그 시리즈",sz=9,c=LG,sp=2)

# ============ P2: 프로젝트 개요 ============
s=sl(); y=M
y=sec(s,y,"프로젝트 개요")
tf=tx(s,M,y,19,16,"",sz=9,c=WHITE)
pa(tf,"WikiEngine은 나무위키, 한국어/영어 위키백과, 뉴스, 웹텍스트 등 6개 소스에서 수집한 12,156,589건의 데이터를 대상으로 한 검색 엔진입니다. 위키 문서를 그대로 사용하지 않고, 실제 커뮤니티 게시판처럼 변환하여 적재했습니다. 위키 [[분류:XXX]]는 태그+카테고리로 변환하고, 뉴스/웹 콘텐츠는 소스별 고정 카테고리를 부여했으며, author_id는 10만 명의 더미 유저에게 균등 배정, created_at은 2020~2025 범위 내 랜덤 생성했습니다.",sz=9,c=WHITE,sp=4)
pa(tf,"",sz=4)
pa(tf,"데이터 소스",sz=10,c=ACCENT,b=True,sp=4)
pa(tf,"  나무위키 (2021.03)           571,364건   JSON   나무마크 본문, 한국어 커뮤니티 문서",sz=8,c=WHITE,sp=3)
pa(tf,"  한국어 위키백과 (2026.03)     739,791건   XML    MediaWiki XML 덤프 (ns=0 일반 문서만)",sz=8,c=WHITE,sp=1)
pa(tf,"  영어 위키백과 (2026.02)     7,139,510건   XML    MediaWiki XML 덤프 (ns=0 일반 문서만)",sz=8,c=WHITE,sp=1)
pa(tf,"  한국어 뉴스                   159,639건   JSON   뉴스 기사 텍스트",sz=8,c=WHITE,sp=1)
pa(tf,"  한국어 웹텍스트             1,284,822건   JSON   웹 크롤링 텍스트",sz=8,c=WHITE,sp=1)
pa(tf,"  C4 한국어 클린              2,261,463건   JSON   한국어 웹 코퍼스 (Common Crawl 정제)",sz=8,c=WHITE,sp=1)
pa(tf,"  합계                       12,156,589건          30개 카테고리, 고유 태그 ~216만 개",sz=8,c=ACCENT,b=True,sp=3)
pa(tf,"",sz=4)
pa(tf,"인프라 구성",sz=10,c=ACCENT,b=True,sp=4)
pa(tf,"  서버 1 (ARM 2코어/12GB): Spring Boot App(Primary) + MySQL(Primary) + Redis(Shard1) + Nginx + Lucene 42GB",sz=8,c=WHITE,sp=3)
pa(tf,"  서버 2 (ARM 2코어/12GB): Spring Boot App(Replica) + MySQL(Replica) + Kafka + Debezium + Redis(Shard2,3)",sz=8,c=WHITE,sp=1)
pa(tf,"  서버 3 (AMD 1GB): Grafana + Loki + InfluxDB(k6)",sz=8,c=WHITE,sp=1)
pa(tf,"  서버 4 (AMD 1GB): Prometheus (500MB, 30일 보존)",sz=8,c=WHITE,sp=1)
pa(tf,"  총 비용: $0/월 (OCI Free Tier)",sz=8,c=GREEN,sp=3)
im(s,IMG/"wiki-search-overview/architecture.png",M,17,w=19)

# ============ P3: 검색 문제 ============
s=sl(); y=M
y=sec(s,y,"1. 검색 엔진 전환 — 정상 상태와 문제 인식")
y=lb(s,y,"1단계: 정상 상태 — MySQL LIKE 검색")
tf=tx(s,M,y,19,4.5,"",sz=9,c=WHITE)
pa(tf,"초기 검색은 MySQL LIKE '%keyword%'로 본문을 검색했습니다. posts 테이블은 id(PK, BIGINT AUTO_INCREMENT), title(VARCHAR 500), content(MEDIUMTEXT, 평균 6,586자), category_id(FK to categories), view_count(BIGINT), like_count(BIGINT), created_at(DATETIME) 구조입니다. categories 테이블은 id(PK), name(VARCHAR UNIQUE), parent_id(BIGINT nullable, 계층 구조)로 posts:categories = N:1 관계입니다. 데이터 규모는 12,156,589건이며, content 컬럼에 위키 마크업이 포함되어 있어 평균 행 크기가 큽니다.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"검색 API: GET /api/v1.0/posts/search?q=keyword&page=0&size=20",sz=8,c=LG,sp=2)
pa(tf,"서버: OCI ARM Ampere A1, 2코어, 12GB RAM, NVMe SSD",sz=8,c=LG,sp=1)
pa(tf,"DB: MySQL 8.0, InnoDB, innodb_buffer_pool_size=4GB",sz=8,c=LG,sp=1)
y+=5

y=lb(s,y,"2단계: 문제 인식 — 타임아웃과 cascade failure",RED)
tf=tx(s,M,y,19,5,"",sz=9,c=WHITE)
pa(tf,"k6 부하 테스트에서 검색 API가 5,000ms 이상 타임아웃되는 현상이 반복적으로 발생했습니다. EXPLAIN으로 확인한 결과 type=ALL, rows=27,443,742 — 전체 테이블을 순차 스캔하고 있었습니다. LIKE '%keyword%'는 와일드카드가 앞에 있어 B-Tree 인덱스를 사용할 수 없고, 27,443,742건의 MEDIUMTEXT를 행마다 문자열 비교해야 합니다.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"더 심각한 문제는 cascade failure였습니다. 검색 쿼리 하나가 HikariCP 커넥션을 수 초간 점유하면, 다른 API(목록 조회, 상세 조회, 게시글 작성)도 커넥션을 얻지 못해 전체 서비스가 마비됩니다. HikariCP maximumPoolSize=10에서 검색 쿼리 3~4개만 동시에 실행되면 나머지 6~7개 커넥션이 모두 대기 상태에 빠집니다.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"1차 시도 — FULLTEXT ngram 인덱스:",sz=9,c=ORANGE,b=True,sp=3)
pa(tf,"12초 -> 6ms (2,100배 개선). 하지만 ngram은 2글자 단위로 토큰을 분리하므로 '대한'이 수백만 건에 매칭되어 posting list 순회가 폭발합니다. 고빈도 토큰에서 다시 5,000ms+ 타임아웃. 또한 12,156,589건 x ngram 토큰 = 인덱스 크기 300GB+ 추정으로, 12GB 서버에서 인덱스 생성 자체가 불가능했습니다.",sz=9,c=WHITE,sp=2)
y+=6

y=lb(s,y,"3단계: 문제 분석 — 왜 MySQL 전문 검색은 한계인가")
tf=tx(s,M,y,19,5,"",sz=9,c=WHITE)
pa(tf,"MySQL FULLTEXT의 구조적 한계 세 가지:",sz=9,c=WHITE,sp=3)
pa(tf,"  1. ngram 토큰 폭발: 2글자 단위 분리로 '대한민국' -> '대한','한민','민국' 3개 토큰. '대한'은 수백만 건에 매칭되어 posting list 순회 비용이 폭발합니다. 한국어 형태소 분석(명사/동사/조사 분리)이 아닌 단순 n-gram이라 false positive가 많습니다.",sz=8,c=WHITE,sp=2)
pa(tf,"  2. 인덱스 크기: 12M건 x 평균 6,586자 x ngram 토큰 = 300GB+ 추정. 서버 디스크 47GB 중 데이터만 30GB를 차지하여 물리적으로 저장 불가.",sz=8,c=WHITE,sp=2)
pa(tf,"  3. boolean mode 한계: 자연어 검색(NL mode)에서 50% 이상 문서에 포함된 term은 자동 제외됩니다(MySQL 공식 문서). 위키 데이터에서 일반적인 한국어 조사/어미가 이 임계값을 넘어 검색 결과에서 빠집니다.",sz=8,c=WHITE,sp=2)
pa(tf,"",sz=2)
pa(tf,"결론: MySQL의 검색 기능으로는 1,215만 건 규모의 한국어 문서를 다룰 수 없습니다. 전문 검색 엔진이 필요합니다.",sz=9,c=ACCENT,b=True,sp=3)
im(s,IMG/"fulltext-ngram-index/like-response-time.png",M,24,w=9)
im(s,IMG/"fulltext-ngram-index/fulltext-response-time.png",10.5,24,w=9)

# ============ P4: Lucene 선택 ============
s=sl(); y=M
y=sec(s,y,"1. 검색 엔진 전환 — 대안 비교와 결과")
y=lb(s,y,"4단계: 대안 비교 — Elasticsearch vs 임베디드 Lucene vs OpenSearch")
tf=tx(s,M,y,19,8,"",sz=9,c=WHITE)
pa(tf,"세 가지 대안을 비교했습니다. 핵심 판단 기준은 '12GB 서버 2대에서 App + MySQL + Redis + Kafka를 모두 운영하면서 검색 엔진을 추가할 수 있는가'였습니다.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=3)
pa(tf,"Elasticsearch:",sz=9,c=RED,b=True,sp=3)
pa(tf,"네이티브 검색 기능(Aggregation, LTR 플러그인, REST API, 동의어 관리)이 풍부합니다. 하지만 단일 노드라도 JVM 2GB+ 메모리가 별도로 필요합니다. HA 구성(3노드)이면 6GB+. 현재 12GB 서버에서 Spring Boot App(JVM 2GB) + MySQL(InnoDB 4GB) + Redis(256MB x 3) + Kafka(JVM 2GB) = 이미 약 9GB를 사용 중입니다. 여기에 ES JVM 2GB를 추가하면 OS + 파일시스템 캐시용 메모리가 부족합니다. 또한 42GB Lucene 인덱스를 ES가 관리하면 ES 프로세스가 추가 메모리를 사용합니다.",sz=8,c=WHITE,sp=2)
pa(tf,"",sz=2)
pa(tf,"임베디드 Lucene (선택):",sz=9,c=GREEN,b=True,sp=3)
pa(tf,"앱 JVM 내에서 동작하여 별도 프로세스가 없습니다. 네트워크 홉 없이 같은 프로세스에서 IndexSearcher를 직접 호출하므로 REST API 직렬화/역직렬화 오버헤드가 없습니다. Nori 한국어 형태소 분석기를 Analyzer 체인에 직접 통합할 수 있고, MMapDirectory로 OS 파일시스템 캐시를 활용합니다. Lucene 10.3.2는 BM25, FeatureField, Rescorer API, SortedSetDocValuesFacets, UnifiedHighlighter 등 ES의 핵심 기능을 대부분 직접 API로 제공합니다.",sz=8,c=WHITE,sp=2)
pa(tf,"한계: 분산 검색이 필요하면 앱 레벨 샤딩을 직접 구현해야 합니다. 현재 2대 서버에서는 rsync로 인덱스를 동기화하고 있으며, 10대+ 확장 시 ES 마이그레이션을 검토합니다.",sz=8,c=ORANGE,sp=2)
pa(tf,"",sz=2)
pa(tf,"OpenSearch Serverless: 관리형이지만 월 $200+. Free Tier 프로젝트에 부적합. 탈락.",sz=8,c=WHITE,sp=2)
y+=9.5

im(s,IMG/"scaleout/step9-load-k6-after-viewcount.png",M,y,w=19)
y+=7

y=lb(s,y,"5단계: 결과 — Lucene 전환 후 수치",GREEN)
tf=tx(s,M,y,19,5,"",sz=9,c=WHITE)
pa(tf,"Lucene 10.3.2 + Nori 형태소 분석기 적용 후:",sz=9,c=WHITE,sp=3)
pa(tf,"  본문 검색: 타임아웃 (5,000ms+) -> 29ms (P95 100ms)",sz=9,c=GREEN,sp=3)
pa(tf,"  BM25 스코어링: title:3, content:1 가중치 적용 (MultiFieldQueryParser)",sz=8,c=WHITE,sp=2)
pa(tf,"  FeatureField 부스팅: viewCount(saturation, w=3.0, pivot=1000) + likeCount(w=2.0, pivot=100)",sz=8,c=WHITE,sp=2)
pa(tf,"  RecencyDecay: 30일 반감기 (DoubleValuesSource)",sz=8,c=WHITE,sp=2)
pa(tf,"  P@10: 0.827 -> 0.853 (+3.2%), PhraseQuery(slop=2) 적용 (15개 테스트 쿼리 수동 평가)",sz=8,c=WHITE,sp=2)
pa(tf,"",sz=3)
pa(tf,"추가 해결한 문제들:",sz=9,c=ACCENT,b=True,sp=3)
pa(tf,"  COUNT(*) 2,038ms 제거: Page -> Slice 전환으로 불필요한 전체 건수 카운트 제거. 에러율 32.53% -> 0%",sz=8,c=WHITE,sp=2)
pa(tf,"  Deferred Join: OFFSET 14,750,000건 페이지네이션이 Lucene보다 38배 느린 문제(2,518ms vs 66ms). Covering Index + Deferred Join 구조로 해결.",sz=8,c=WHITE,sp=2)

# ============ P5: 캐시 + 자동완성 ============
s=sl(); y=M
y=sec(s,y,"2. 캐시 전략 + 자동완성")
y=lb(s,y,"정상 상태와 문제")
tf=tx(s,M,y,19,4,"",sz=9,c=WHITE)
pa(tf,"Lucene 검색 자체는 29ms로 빨라졌지만, 동일한 검색어에 대해 매번 Lucene 인덱스를 탐색하는 것은 CPU 낭비입니다. 검색 트래픽은 Zipf 분포를 따라 소수의 인기 검색어가 전체 트래픽의 대부분을 차지합니다. k6 부하 테스트에서 BM25 + FeatureField + RecencyDecay 스코어링이 CPU를 지속적으로 사용하여, 100 VU에서 CPU 80~100% 포화가 발생했습니다.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"Redis vs CDN vs Caffeine 비교:",sz=9,c=ACCENT,b=True,sp=3)
pa(tf,"CDN은 이미지/정적 리소스에 적합하지만, 검색 결과는 키워드+페이지 조합으로 캐시 키가 폭발합니다. Redis(ElastiCache)는 월 3만원+이지만, Free Tier 서버 내 Redis는 추가 비용 없습니다. Caffeine은 JVM 내 L1 캐시로 네트워크 지연 없이 0.1ms 미만 응답이 가능합니다. 단일 서버에서는 Caffeine만으로 충분하지만, 스케일아웃 시 인스턴스 간 캐시 불일치가 발생합니다.",sz=8,c=WHITE,sp=2)
y+=4.8

y=lb(s,y,"결과 — Caffeine L1 + Redis L2 2계층 캐시",GREEN)
tf=tx(s,M,y,19,2.5,"",sz=9,c=WHITE)
pa(tf,"  Caffeine L1: 73% 히트 (같은 인스턴스 반복 요청, 0.1ms)",sz=9,c=GREEN,sp=3)
pa(tf,"  Redis L2: 9% 추가 히트 (다른 인스턴스 캐시 미스 커버, P95 2.5ms)",sz=9,c=GREEN,sp=2)
pa(tf,"  합산: 82% 히트율, Origin(Lucene/MySQL) 도달률 19%",sz=9,c=GREEN,sp=2)
pa(tf,"  전체 응답: 776ms -> 54ms (14.4배 개선)",sz=9,c=GREEN,sp=2)
y+=3

im(s,IMG/"redis-l2-cache/after-03-spring-boot-system-tiered.png",M,y,w=19)
y+=7

y=lb(s,y,"자동완성 — Trie에서 Redis flat KV로")
tf=tx(s,M,y,19,5,"",sz=9,c=WHITE)
pa(tf,"자동완성을 위해 먼저 Trie 자료구조를 구현했습니다. 검색 로그에서 인기 검색어를 추출하고 Trie에 적재하여, 접두사 입력 시 DFS로 Top-10 제안을 반환합니다. 한국어 자모 분해(초성/중성/종성)를 적용하여 '삼ㅅ' 입력 시 '삼성전자'가 제안되도록 구현했습니다.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"하지만 Trie는 JVM 힙에 존재하므로 App 2대에서 각각 독립적인 Trie를 가져야 하고, 인스턴스 간 일관성 문제가 발생합니다. 또한 DFS 순회 + Copy-on-Write가 CPU/메모리 부담입니다. Redis flat KV로 전환하여 접두사를 키로, Top-10 제안을 JSON 배열로 사전 계산합니다. Redis GET O(1) 단일 명령으로 조회하므로 DFS 순회가 없고, 공유 저장소이므로 어떤 App이 응답해도 동일한 결과입니다.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"이후 @Scheduled 배치를 Spring Batch Job+Tasklet으로 전환하여 JobRepository에 실행 이력이 자동 기록되고, FAILED 상태에서 재시작이 가능해졌습니다.",sz=8,c=ORANGE,sp=3)

# ============ P6: 분산 - 부하테스트 ============
s=sl(); y=M
y=sec(s,y,"3. 분산 아키텍처 — 단일 서버 한계 발견")
y=lb(s,y,"정상 상태와 문제 — k6 stress 테스트 (200 VU, 25분)")
tf=tx(s,M,y,19,5,"",sz=9,c=WHITE)
pa(tf,"Lucene + Caffeine 캐시까지 적용한 상태에서 k6 stress 테스트(200 VU, 25분)를 수행했습니다. ARM 2코어 서버에서 CPU가 100% 포화되었습니다. Lucene BM25 + Nori 형태소 분석이 CPU-intensive하여, Caffeine 캐시 미스(18%)가 발생할 때마다 CPU가 급격히 올라갑니다. Load Average가 20+까지 치솟으며 모든 Tomcat 스레드가 점유됩니다.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"  P95: 1,413ms (SLA 300ms 초과)",sz=9,c=RED,sp=3)
pa(tf,"  에러율: 13.25% (100 VU에서)",sz=9,c=RED,sp=2)
pa(tf,"  처리량: ~30 req/s",sz=9,c=RED,sp=2)
pa(tf,"  CPU: 100% 지속 포화",sz=9,c=RED,sp=2)
pa(tf,"  JVM G1GC, HikariCP 튜닝 시도했으나 근본 해결 안 됨 (CPU가 병목)",sz=8,c=ORANGE,sp=3)
y+=5.5

im(s,IMG/"stress-test-tuning/A-2-k6-grafana-overview.png",M,y,w=19)
y+=7.5

y=lb(s,y,"분석 — 왜 튜닝으로 해결할 수 없는가")
tf=tx(s,M,y,19,4,"",sz=9,c=WHITE)
pa(tf,"문제의 근본 원인은 2코어 ARM CPU의 물리적 한계입니다. BM25 스코어링은 term마다 posting list를 순회하면서 TF/IDF를 계산하고, Nori는 한국어 문장을 형태소 단위로 분해합니다. 이 두 연산이 CPU-bound입니다. 캐시 히트율을 82%까지 올려도 나머지 18%의 Origin 요청이 CPU를 포화시킵니다.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"JVM 튜닝(G1GC 힙/스레드, HikariCP 풀 사이즈)은 메모리/IO 병목에 효과적이지만, CPU 병목에는 소용없습니다. 물리적으로 코어를 추가하는 수평 확장(스케일아웃)이 필요합니다. 스케일아웃의 전제조건: Caffeine 캐시, Trie 자동완성, TokenBlacklist가 모두 JVM 힙에 있어 인스턴스 간 상태가 공유되지 않습니다. 먼저 Stateless로 전환해야 합니다.",sz=9,c=WHITE,sp=3)

# ============ P7: 분산 컴포넌트 ============
s=sl(); y=M
y=sec(s,y,"3. 분산 아키텍처 — 6개 컴포넌트 도입")
tf=tx(s,M,y,19,24,"",sz=9,c=WHITE)
pa(tf,"1. Stateless 전환 + Redis L2 캐시",sz=10,c=ACCENT,b=True,sp=4)
pa(tf,"Caffeine(L1 캐시), Trie(자동완성), TokenBlacklist(JWT 무효화)를 모두 Redis로 외부화했습니다. 이제 App 인스턴스에 상태가 없으므로 자유롭게 추가/제거할 수 있습니다. Trie는 Redis flat KV O(1)로 전환(앞서 설명), TokenBlacklist는 Redis Set + TTL로 이관했습니다.",sz=8,c=WHITE,sp=2)
pa(tf,"",sz=3)
pa(tf,"2. MySQL GTID Replication (R/W 분리)",sz=10,c=ACCENT,b=True,sp=4)
pa(tf,"Primary(쓰기) + Replica(읽기) 구성. AbstractRoutingDataSource + LazyConnectionDataSourceProxy로 @Transactional(readOnly) 어노테이션 기반 자동 라우팅합니다. HikariCP를 Primary 5 + Replica 15 커넥션풀로 분리했습니다. 검색은 대부분 읽기이므로 Replica가 70%+ 부하를 흡수합니다. Lag은 0~1초로 최종 일관성(Eventual Consistency) 허용.",sz=8,c=WHITE,sp=2)
pa(tf,"",sz=3)
pa(tf,"3. Nginx L7 least_conn + App 2대",sz=10,c=ACCENT,b=True,sp=4)
pa(tf,"App 2대로 확장하고 Nginx로 L7 로드 밸런싱합니다. least_conn 알고리즘으로 현재 활성 연결이 적은 서버에 요청을 보냅니다. Lucene 인덱스는 rsync + cron으로 Primary -> Replica 동기화합니다.",sz=8,c=WHITE,sp=2)
pa(tf,"결과: 에러율 13.25% -> 0%, P95 2,300ms -> 158ms, 평균 482ms -> 37ms",sz=8,c=GREEN,sp=2)
pa(tf,"",sz=3)
pa(tf,"4. Redis INCR + Write-Behind (조회수 동시성)",sz=10,c=ACCENT,b=True,sp=4)
pa(tf,"게시글 상세 조회(GET) API에서 viewCount UPDATE가 실행되는데, R/W 분리 후 읽기 트랜잭션에서 UPDATE를 시도하여 에러가 발생합니다(Read-only connection). Redis INCR로 원자적 카운팅 후 30초 배치로 DB에 플러시하는 Write-Behind 패턴으로 해결했습니다. GET 요청에서 쓰기가 완전히 제거되어 에러율 11.10% -> 0%.",sz=8,c=WHITE,sp=2)
pa(tf,"",sz=3)
pa(tf,"5. Debezium + Kafka CDC",sz=10,c=ACCENT,b=True,sp=4)
pa(tf,"PostService에서 MySQL UPDATE와 Lucene updateDocument를 직접 호출하는 dual-write가 강한 결합과 성능 문제를 일으켰습니다(게시글 생성 5,315ms). Spring Event -> @ApplicationModuleListener -> Debezium+Kafka CDC 3단계로 진화시켰습니다. MySQL binlog를 Debezium이 캡처하여 Kafka topic으로 발행하고, Consumer가 Lucene 인덱스와 L1/L2 캐시를 자동 업데이트합니다. 게시글 생성 5,315ms -> 33ms (160배). DLQ + AckMode.RECORD로 에러 핸들링을 강화하여, 실패 메시지는 Dead Letter Topic에 격리되고 9회 재시도 후에도 실패하면 사후 분석이 가능합니다.",sz=8,c=WHITE,sp=2)
pa(tf,"",sz=3)
pa(tf,"6. Redis 3노드 Consistent Hashing",sz=10,c=ACCENT,b=True,sp=4)
pa(tf,"단일 Redis에서 KEYS 명령이 34.6ms SLOWLOG에 기록되었습니다. 자동완성 배치(buildPrefixTopK)가 매시간 수만 개 키를 동시에 SET하면서, 실시간 GET/INCR과 같은 싱글스레드에서 경합합니다. KEYS -> SCAN으로 전환하고, 3노드 Consistent Hashing(MurmurHash3)으로 자동완성/캐시/블랙리스트 워크로드를 물리적으로 격리했습니다.",sz=8,c=WHITE,sp=2)

# ============ P8: 분산 결과 ============
s=sl(); y=M
y=sec(s,y,"3. 분산 아키텍처 — 검증 결과")
y=lb(s,y,"k6 stress 테스트 (200 VU, 분산 아키텍처)",GREEN)
tf=tx(s,M,y,19,4,"",sz=9,c=WHITE)
pa(tf,"분산 아키텍처(2 App + MySQL Replication + Redis 3샤드 + Kafka CDC) 전환 후 200 VU stress 테스트를 수행했습니다:",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"  에러율: 13.25% (단일) -> 0.09% (분산)     에러 거의 제거",sz=9,c=GREEN,sp=3)
pa(tf,"  처리량: 30 req/s -> 109 req/s              3.6배 증가",sz=9,c=GREEN,sp=2)
pa(tf,"  P95: 2,300ms -> 190ms                      12배 개선",sz=9,c=GREEN,sp=2)
pa(tf,"  100 VU P95: 200ms                          SLA(300ms) 충족",sz=9,c=GREEN,sp=2)
pa(tf,"",sz=2)
pa(tf,"MySQL(InnoDB 버퍼풀 히트율 100%, Row Lock 0), Redis(OPS 정상), Kafka(Consumer Lag 0~3K -> 0 수렴, CDC Lag 피크 40ms), Nginx(Active Conn 200+) 모두 여유가 있었고, App CPU만이 80-100%로 여전히 근본 병목임을 소거법으로 확인했습니다.",sz=8,c=WHITE,sp=3)
pa(tf,"Replication Lag 0~1초 안정. Primary SELECT ~10 ops/s, Replica SELECT ~70 ops/s. 수신 데이터 242 MiB.",sz=8,c=LG,sp=2)
y+=5.5

im(s,IMG/"distributed-stability/phase16-stress-k6-overview.png",M,y,w=19)
y+=7.5

im(s,IMG/"cdc/phase14-cdc-k6-overview.png",M,y,w=19)

# ============ P9: 검색 품질 - 동의어/오타/snippet ============
s=sl(); y=M
y=sec(s,y,"4. 검색 품질 — 동의어 확장 + 오타 교정 + Snippet")
y=lb(s,y,"문제 인식 — 세 가지 검색 품질 한계")
tf=tx(s,M,y,19,4,"",sz=9,c=WHITE)
pa(tf,'1. 동의어 미지원 (Recall 손실): "AI"를 검색하면 "인공지능" 문서가 누락됩니다. Lucene은 정확한 term 매칭만 하므로, 약어/동의어 관계를 알 수 없습니다. "DB" -> "데이터베이스", "ML" -> "머신러닝"도 마찬가지입니다.',sz=9,c=WHITE,sp=3)
pa(tf,'2. 오타 교정 미지원 (검색 실패): "프로그래링" 검색 시 결과 0건. "컴퓨텨" 검색 시 결과 0건. 검색 결과 0건은 사용자 이탈의 주요 원인입니다.',sz=9,c=WHITE,sp=3)
pa(tf,"3. Snippet 무관: 문서 앞 150자를 잘라서 snippet으로 표시. '가비지 컬렉션' 검색 시 '자바는 제임스 고슬링이 1995년에...' (GC 언급 없음). 검색어가 본문 중간에 있으면 snippet에 나타나지 않습니다.",sz=9,c=WHITE,sp=3)
y+=4.5

y=lb(s,y,"대안 비교 — 동의어 처리")
tf=tx(s,M,y,19,3,"",sz=9,c=WHITE)
pa(tf,"인덱스 타임 동의어(SynonymGraphFilter): 인덱스에 동의어 term을 추가하여 document frequency를 인위적으로 높입니다. 'AI'를 인덱싱할 때 '인공지능'도 함께 추가하면, '인공지능'의 DF가 실제보다 부풀려져 BM25 IDF 계산이 왜곡됩니다. 해당 term의 가중치가 낮아져 오히려 검색 품질이 저하될 수 있습니다. (OpenSource Connections 'Solr Synonyms Mea Culpa'에서 실사례로 경고)",sz=8,c=WHITE,sp=3)
pa(tf,"쿼리 타임 DB 기반 확장 (선택): 인덱스 term 통계가 불변이므로 IDF 왜곡이 없습니다. 동의어 추가/삭제 시 재색인이 불필요하고, DB에서 즉시 반영됩니다. Caffeine 캐시(TTL 30분)로 DB 조회 비용을 완화합니다.",sz=8,c=GREEN,sp=3)
y+=3.5

y=lb(s,y,"결과",GREEN)
im(s,IMG/"search-query-enhancement/phase18-before-search-ai.png",M,y,w=9)
im(s,IMG/"search-query-enhancement/phase18-after-search-ai.png",10.5,y,w=9)
y+=6
tf=tx(s,M,y,19,4,"",sz=9,c=WHITE)
pa(tf,"  'AI' 검색: 영문 'Ai' 문서만 -> '인공지능' 1위 (동의어 확장)",sz=9,c=GREEN,sp=3)
pa(tf,"  '프로그래링' -> '혹시 프로그래밍을 찾으셨나요?' (DirectSpellChecker, Damerau-Levenshtein)",sz=9,c=GREEN,sp=2)
pa(tf,"  Snippet: UnifiedHighlighter + snippetSource 500자 StoredField -> 검색어 주변 맥락 추출",sz=9,c=GREEN,sp=2)
pa(tf,"  Nori 사용자 사전 158,539개 (open-korean-text Apache 2.0) 적용으로 복합어 보존",sz=8,c=WHITE,sp=2)
pa(tf,"  12,156,589건 무중단 재색인 ~2시간 (Directory Swap + SearcherManager 재생성)",sz=8,c=WHITE,sp=2)

# ============ P10: LTR ============
s=sl(); y=M
y=sec(s,y,"5. LTR 재랭킹 — XGBoost LambdaMART")
y=lb(s,y,"문제 — 수동 가중치의 한계")
tf=tx(s,M,y,19,3,"",sz=9,c=WHITE)
pa(tf,'"자바" 검색 시 "자바 더 헛(스타워즈)"이 1위, "자바(프로그래밍 언어)"가 4위로 밀립니다. viewCount가 모두 0~1로 동일하여 BM25 텍스트 스코어만으로 순위가 결정됩니다. Linear Model은 ES boosts와 동일한 선형 회귀이므로 같은 피처로 학습해도 수동 가중치와 거의 동일한 결과가 나옵니다 (OpenSource Connections 분석).',sz=9,c=WHITE,sp=3)
y+=3.2

y=lb(s,y,"해결 — LambdaMART + LLM-as-a-Judge")
tf=tx(s,M,y,19,6,"",sz=9,c=WHITE)
pa(tf,"XGBoost LambdaMART (tree model)는 피처 간 interaction을 학습할 수 있습니다. titleLength가 짧으면서 tagOverlap이 높으면 동음이의어 분기 문서라는 비선형 관계를 포착합니다.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"학습 데이터: 사용자 트래픽이 거의 없어 클릭 로그가 없으므로, LLM-as-a-Judge(Gemini)로 relevance 판정을 대신합니다. SIGIR 2024 Thomas et al. — GPT-4의 relevance 판정이 crowdsource annotator와 Cohen Kappa 동등 수준. 3회 호출 평균 반올림으로 LLM의 비결정성을 대응합니다.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"1차 실행 실패 (98%): 라운드 간 딜레이 2초로 분당 30요청 -> Gemini 15 RPM 초과. Spring AI가 HTTP 429를 NonTransientAiException으로 분류하여 재시도하지 않았고, ArrayList 메모리 저장이라 실패 미인지. 5초 딜레이 + 지수 백오프 + CSV 디스크 저장 + resume으로 해결.",sz=8,c=ORANGE,sp=3)
pa(tf,"",sz=2)
pa(tf,"추론: XGBoost4J (ARM64 네이티브 JAR 번들, inplace_predict thread-safe). Rescorer Top-200 재랭킹.",sz=8,c=WHITE,sp=2)
pa(tf,"14개 피처: bm25Title, bm25Content, queryTermCoverage(title/content), exactTitleMatch, titleLength, contentLength, freshnessDays, viewCount, likeCount, tagOverlap, categoryId, queryLength, bm25Snippet",sz=8,c=LG,sp=2)
y+=7

y=lb(s,y,"결과",GREEN)
tf=tx(s,M,y,19,2,"",sz=9,c=WHITE)
pa(tf,"  NDCG@10: 0.6910 -> 0.7387 (+4.8%p, 5-Fold CV)    train 0.9228은 과적합 — CV가 정직한 수치",sz=9,c=GREEN,sp=3)
pa(tf,'  "자바" 1위: 자바(프로그래밍 언어)    Feature Importance: titleLength(1.0) > tagOverlap(0.9) > bm25Title(0.8)',sz=9,c=GREEN,sp=2)
pa(tf,"  한계: 2코어 ARM LTR ON -> 72배 악화 (42ms->3,088ms). CPU 완전 포화. LTR_ENABLED=false로 비활성화.",sz=9,c=RED,sp=3)
y+=2.5

im(s,IMG/"search-ltr-ranking/phase19-ltr-before-java-frontend.png",M,y,w=9)
im(s,IMG/"search-ltr-ranking/phase19-ltr-after-java-frontend.png",10.5,y,w=9)
y+=5.5
im(s,IMG/"search-ltr-ranking/phase19-ltr-training-result.png",M,y,w=19)

# ============ P11: Facet + 카테고리 ============
s=sl(); y=M
y=sec(s,y,"6. 카테고리 자동 분류 + Facet 네이티브 전환")
y=lb(s,y,"문제 — namespace 카테고리로는 Facet이 무의미")
tf=tx(s,M,y,19,2.5,"",sz=9,c=WHITE)
pa(tf,"위키 데이터 임포트 시 할당된 카테고리가 '일반 문서' 97%로 편중되어 있어 Facet을 표시해도 의미가 없었습니다. 또한 카테고리 필터링은 Occur.FILTER 절로 구현했지만, Facet 집계는 DB GROUP BY로 Top-1,000건의 근사 집계만 가능했습니다. SortedSetDocValuesFacetField 추가 + 전체 재색인이 필요했습니다.",sz=9,c=WHITE,sp=3)
y+=3

y=lb(s,y,"해결 — 28개 주제별 자동 분류 + 1회 재색인으로 통합",GREEN)
tf=tx(s,M,y,19,5,"",sz=9,c=WHITE)
pa(tf,"키워드 기반 배치 분류 서비스(CategoryClassificationService)로 28개 주제별 카테고리(컴퓨터 과학, 수학, 물리학, 역사, 음악, 게임, 스포츠 등)를 자동 분류했습니다. 90건 수동 검증 결과 정확도 ~83% (애매 항목 제외 시 92%). 키워드 간 우선순위/배타성 로직이 없어 경계 케이스에서 정확도 저하가 있지만, Facet 탐색용으로 충분합니다.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"1회 재색인으로 통합 반영:",sz=9,c=ACCENT,b=True,sp=3)
pa(tf,"  SortedSetDocValuesFacetField('category', categoryName) — Lucene 네이티브 Facet",sz=8,c=WHITE,sp=2)
pa(tf,"  TextField('tags', tagNames, Store.YES) — 216만 고유 태그 검색 가능 (Facet은 고카디널리티로 제거)",sz=8,c=WHITE,sp=2)
pa(tf,"  Nori 사용자 사전 158,539개, snippetSource clean text",sz=8,c=WHITE,sp=2)
pa(tf,"  RefreshListener 기반 SortedSetDocValuesReaderState 캐싱 (LUCENE-7905 OrdinalMap 빌드 비용 대응)",sz=8,c=WHITE,sp=2)
pa(tf,"  MultiCollectorManager로 TopDocs + FacetsCollector 단일 패스 수집 (searcher.search 2회 -> 1회)",sz=8,c=WHITE,sp=2)
y+=6

im(s,IMG/"search-ltr-ranking/phase19-after-category-facets-distribution.png",M,y,w=19)

# ============ P12: 콘텐츠 필터링 ============
s=sl(); y=M
y=sec(s,y,"7. 콘텐츠 필터링 — 운영 안전장치")
y=lb(s,y,"문제 — 게시글 작성 API가 열려있고 필터가 없다")
tf=tx(s,M,y,19,2.5,"",sz=9,c=WHITE)
pa(tf,"POST /api/v1.0/posts로 아무 콘텐츠나 저장 가능합니다. 유해 콘텐츠가 검색 결과와 자동완성에 그대로 노출됩니다. k6 부하 테스트에서도 임의 문자열이 게시글로 생성되어 검색 품질을 오염시킵니다.",sz=9,c=WHITE,sp=3)
y+=3

y=lb(s,y,"해결 — Aho-Corasick + 블라인드 + Negative Caching",GREEN)
tf=tx(s,M,y,19,7,"",sz=9,c=WHITE)
pa(tf,"금칙어 필터링: String.contains() 루프는 O(N*M)으로 금칙어 16,090개면 매 텍스트마다 16,090번 순회합니다. Aho-Corasick은 Trie + failure link로 텍스트를 한 번만 순회하면서 모든 패턴을 동시 매칭합니다. O(N+Z)로 금칙어 수에 무관합니다. 한국어 Trie(부분 일치)와 영어 Trie(단어 경계, Scunthorpe 방지)를 분리하여 운영합니다.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"블라인드: Occur.MUST_NOT으로 검색에서 제외합니다. 인덱스에서 삭제하면 복원 시 재인덱싱이 필요하지만, blinded=true 필드를 두고 MUST_NOT으로 필터하면 blinded=false 업데이트만으로 즉시 복원됩니다.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"Negative Caching: 앱 기동 직후 Lucene 인덱스 로딩 전에 검색하면 0건 결과가 캐시됩니다. 30초 TTL로 빈 결과를 캐시하되 빠르게 만료시킵니다. 빈 결과를 아예 캐시하지 않으면 cache penetration(동일 쿼리가 매번 origin까지 관통)이 발생합니다. RFC 2308, AWS CloudFront(negative TTL 기본값 5초) 참고.",sz=9,c=WHITE,sp=3)
pa(tf,"",sz=2)
pa(tf,"자동완성 Lucene fallback: Nori-analyzed title 필드로 PrefixQuery를 하면 '성매' -> Nori -> '성'으로 분해 -> '성악가' 등 무관한 결과. title_raw StringField(untokenized)를 추가하여 원본 prefix 매칭.",sz=8,c=WHITE,sp=3)
y+=8

im(s,IMG/"search-content-filter/phase20-autocomplete-normal-java.png",M,y,w=9)
im(s,IMG/"search-content-filter/phase20-blind-after-test-search.png",10.5,y,w=9)

# ============ P13: RAG ============
s=sl(); y=M
y=sec(s,y,"8. AI 검색 요약 — RAG 파이프라인")
y=lb(s,y,"문제 — 문서 목록만 제공, 직접 답변 불가")
tf=tx(s,M,y,19,2.5,"",sz=9,c=WHITE)
pa(tf,"기존 검색 결과는 10건의 게시글 링크(제목 + snippet)뿐입니다. 사용자가 직접 2~3개 게시글을 클릭해서 읽어야 원하는 답을 찾을 수 있습니다. 기존 AI 요약 구현은 검색 결과와 무관하게 LLM에 쿼리만 전달하여 할루시네이션 위험이 높고, 출처 인용이 불가능했습니다.",sz=9,c=WHITE,sp=3)
y+=3

y=lb(s,y,"해결 — BM25 Top-5 -> LLM Context -> SSE 스트리밍",GREEN)
tf=tx(s,M,y,19,8,"",sz=9,c=WHITE)
pa(tf,"RAG 파이프라인 5단계:",sz=9,c=ACCENT,b=True,sp=3)
pa(tf,"  1. Query Understanding — 오타 교정 + 동의어 확장 (이전 구현 재활용)",sz=8,c=WHITE,sp=2)
pa(tf,"  2. Retrieval — Lucene BM25 Top-5 문서 검색",sz=8,c=WHITE,sp=1)
pa(tf,"  3. Context 구성 — 5개 문서의 500자씩 프롬프트에 주입 (~2,000 토큰)",sz=8,c=WHITE,sp=1)
pa(tf,"  4. Generation — Spring AI + Gemini 2.0 Flash, SseEmitter + Virtual Thread SSE 스트리밍",sz=8,c=WHITE,sp=1)
pa(tf,"  5. Response — AI 답변 + [문서 N] 인라인 출처 배지 파싱 -> 게시글 링크",sz=8,c=WHITE,sp=1)
pa(tf,"",sz=3)
pa(tf,"왜 BM25가 Dense Retrieval보다 적합한가:",sz=9,c=ACCENT,b=True,sp=3)
pa(tf,"위키피디아 기술 용어 키워드 검색이 주 패턴. 키워드 전용 쿼리에서 BM25 NDCG 0.88 > Dense 0.65. 'AI'->'인공지능' 수준은 동의어 확장으로 이미 해결. 벡터 검색은 '환경 오염'->'대기질' 같은 의미 검색이 실측 확인되면 Lucene KnnFloatVectorField(HNSW) + Reciprocal Rank Fusion으로 전환 예정.",sz=8,c=WHITE,sp=2)
pa(tf,"",sz=3)
pa(tf,"할루시네이션 방지: 시스템 프롬프트에 '제공된 문서만 참고' + [문서 N] 인용 강제 + BM25 스코어 임계값 미만 시 AI 스킵 + 네비게이션 의도(네이버/구글) 스킵.",sz=8,c=WHITE,sp=3)
pa(tf,"",sz=3)
pa(tf,"비용 최적화:",sz=9,c=ACCENT,b=True,sp=3)
pa(tf,"  Rate Limiting: Redis Lua 스크립트 atomic (INCR+EXPIRE non-atomic 문제 해결), 10 RPM 전역 공유",sz=8,c=WHITE,sp=2)
pa(tf,"  동일 쿼리 캐싱: Redis TTL 30분 -> LLM 비용 40~60% 절감",sz=8,c=GREEN,sp=1)
pa(tf,"  Gemini 무료 티어 활용: 15 RPM, 1,000 RPD. 월 $0",sz=8,c=WHITE,sp=1)
pa(tf,"  Grafana 7패널 대시보드: RPM, 응답시간, 토큰, 피드백, 일별 비용 추정",sz=8,c=WHITE,sp=1)
y+=9.5

im(s,IMG/"search-rag/ai-summary-full-with-results.png",M,y,w=9)
im(s,IMG/"search-rag/grafana-llm-panels.png",10.5,y,w=9)

# ============ P14: 실패와 교훈 ============
s=sl(); y=M
y=sec(s,y,"9. 실패와 교훈")
tf=tx(s,M,y,19,25,"",sz=9,c=WHITE)
pa(tf,"1. LTR ON 시 72배 성능 악화",sz=10,c=ORANGE,b=True,sp=4)
pa(tf,"원인: Rescore window 200에서 문서당 14개 피처 추출이 CPU-intensive합니다. BM25 3필드 x 200문서 = 600회 Scorer 생성 + Nori 토큰화 200문서 x 3회 = 600회+ 형태소 분석. 100 VU 동시 요청 시 2코어 ARM에서 CPU 완전 포화 -> queueing avalanche -> LTR 무관 API(자동완성, 목록 조회)까지 동반 악화. App CPU 피크 160% (2코어 중 80%+80%), Load Average 60.",sz=8,c=WHITE,sp=2)
pa(tf,"교훈: 기능 검증과 프로덕션 적용은 별개입니다. 현업에서는 피처를 인덱스 타임에 사전 계산(stored field)하거나, 피처 캐싱, 전용 다코어 서버에서 LTR을 처리합니다. LTR_ENABLED=false로 비활성화하고 인프라 확장 시 재활성화 예정.",sz=8,c=WHITE,sp=2)
pa(tf,"",sz=4)
pa(tf,"2. LLM 학습 데이터 생성 98% 실패",sz=10,c=ORANGE,b=True,sp=4)
pa(tf,"원인: 라운드 간 딜레이 2초로 3회 호출이 ~6초, 분당 30요청으로 Gemini 15 RPM 초과. Spring AI 기본 retry가 HTTP 429를 NonTransientAiException(4xx)으로 분류하여 재시도하지 않았습니다. 데이터가 메모리 전용(ArrayList)이라 status API에 성공/실패 구분이 없어 98% 실패를 인지하지 못했습니다.",sz=8,c=WHITE,sp=2)
pa(tf,"교훈: 외부 API rate limit을 정확히 계산(15 RPM -> 5초 간격 = 12 RPM). 데이터는 즉시 디스크 저장(CSV append + flush). status API에 성공/실패 분리 필수. resume 기능(완료된 건 건너뛰기)으로 crash-safe.",sz=8,c=WHITE,sp=2)
pa(tf,"",sz=4)
pa(tf,"3. CDC 배포 후 검색 미반영",sz=10,c=ORANGE,b=True,sp=4)
pa(tf,"원인: 멀티 인스턴스 환경에서 CDC Consumer와 Lucene IndexWriter의 위치 불일치. App 2(Kafka와 같은 서버)에서 CDC Consumer가 이벤트를 수신했지만, App 2는 Lucene Replica(IndexWriter가 null)라 인덱싱이 skip. App 1은 docker-compose에서 SPRING_KAFKA_BOOTSTRAP_SERVERS 환경변수 매핑이 빠져 Kafka 연결 자체가 실패.",sz=8,c=WHITE,sp=2)
pa(tf,"교훈: 임베디드 검색 인덱스는 분산 환경에서 Consumer-Writer 위치 보장이 어렵습니다. CDC 파이프라인은 end-to-end 검증 테스트가 필수입니다. 이 경험에서 현업에서 검색 인덱스가 앱에 embedded 되지 않는 이유를 체감했습니다.",sz=8,c=WHITE,sp=2)
pa(tf,"",sz=4)
pa(tf,"4. Flyway 마이그레이션 Replica 미전파",sz=10,c=ORANGE,b=True,sp=4)
pa(tf,"원인: MySQL Replication이 끊겨 DDL이 Replica에 전파되지 않았습니다. ddl-auto:validate(운영)와 update(로컬)의 차이로, 로컬에서는 자동 처리되던 스키마 변경이 운영에서 실패했습니다. ai_summary_feedback 테이블과 blinded 컬럼이 Replica에 없어 읽기 쿼리가 전부 실패.",sz=8,c=WHITE,sp=2)
pa(tf,"교훈: 배포 전 Replication 상태 + Flyway 마이그레이션 체크리스트 필수. docker-compose environment 블록에 새 환경변수 추가 누락 방지.",sz=8,c=WHITE,sp=2)
pa(tf,"",sz=4)
pa(tf,"5. snippetSource에 raw 위키 마크업 저장",sz=10,c=ORANGE,b=True,sp=4)
pa(tf,"원인: 위키 마크업([include(틀:XXX)], {{생물 분류}} 등)이 stored field에 그대로 저장. UnifiedHighlighter가 마크업 토큰에서 offset 불일치 -> 빈 snippet 반환.",sz=8,c=WHITE,sp=2)
pa(tf,"교훈: raw 데이터를 stored field에 저장하는 것은 안티패턴입니다. Wikipedia CirrusSearch 패턴: source_text(raw)와 text(clean)을 별도 필드로 관리. 인덱스 타임에 마크업을 정제한 clean text를 저장합니다.",sz=8,c=WHITE,sp=2)

# ============ P15: 최종 아키텍처 + 수치 총정리 ============
s=sl(); y=M
y=sec(s,y,"최종 아키텍처 + 핵심 수치")

im(s,"/tmp/wikiengine-architecture.png",M,y,w=19)
y+=11

tf=tx(s,M,y,19,16,"",sz=9,c=WHITE)
pa(tf,"검색 성능 (k6 부하 테스트 실측)",sz=9,c=ACCENT,b=True,sp=3)
pa(tf,"                                 Before                       After                        측정 조건",sz=7,c=LG,sp=3)
pa(tf,"  본문 검색            LIKE 타임아웃 (5,000ms+)     Lucene BM25 29ms (P95 100ms)      EXPLAIN rows=27,443,742",sz=8,c=WHITE,sp=2)
pa(tf,"  B-Tree 자동완성      LIKE 타임아웃 (5,000ms+)     8ms (rows 27.44M -> 1건 스캔)     복합 인덱스 idx_title_viewcount",sz=8,c=WHITE,sp=1)
pa(tf,"  FULLTEXT ngram       12,766ms                     6ms (2,100배)                      57만건, 인덱스 6.7GB",sz=8,c=WHITE,sp=1)
pa(tf,"  Redis 자동완성       Trie DFS 5ms                 Redis GET 11ms (P95 68ms)          flat KV O(1), 5,000 prefix 키",sz=8,c=WHITE,sp=1)
pa(tf,"  캐시 전체 응답       776ms                        54ms (14.4배)                       Caffeine 99.9% 히트",sz=8,c=WHITE,sp=1)
pa(tf,"  L1+L2 히트율         0% (캐시 없음)               82% (L1 73% + L2 9%)               Origin 도달률 19%",sz=8,c=WHITE,sp=1)
pa(tf,"  Redis 메모리         -                            73MB / 256MB (28.4%)               Eviction 0, Lettuce P95 2.5ms",sz=8,c=WHITE,sp=1)
pa(tf,"",sz=3)
pa(tf,"분산 아키텍처 (k6 stress 200 VU, 25분)",sz=9,c=ACCENT,b=True,sp=3)
pa(tf,"                          단일 서버                    분산 (2 App)                     변화",sz=7,c=LG,sp=3)
pa(tf,"  에러율 (쿼리)        32.53% (COUNT* 타임아웃)      0%                                Page -> Slice 전환",sz=8,c=WHITE,sp=2)
pa(tf,"  에러율 (부하 100VU)  13.25%                       0%                                스케일아웃 효과",sz=8,c=WHITE,sp=1)
pa(tf,"  에러율 (부하 200VU)  -                            0.09%                              MySQL/Redis/Kafka 여유",sz=8,c=WHITE,sp=1)
pa(tf,"  평균 응답 (100VU)    482ms                        37ms (92% 감소)                    Nginx least_conn + App 2대",sz=8,c=WHITE,sp=1)
pa(tf,"  P95 (100VU)          2,300ms                      190ms (12배)                       SLA 300ms 충족",sz=8,c=WHITE,sp=1)
pa(tf,"  처리량               ~30 req/s                    109 req/s (3.6배)                   피크 기준",sz=8,c=WHITE,sp=1)
pa(tf,"  상세조회 에러        11.10% (DB UPDATE 충돌)       0% (Redis INCR)                    Write-Behind 30초 배치",sz=8,c=WHITE,sp=1)
pa(tf,"  게시글 생성          5,315ms (dual-write)         33ms (CDC, 160배)                   Debezium+Kafka",sz=8,c=WHITE,sp=1)
pa(tf,"  KEYS 블로킹          34.6ms SLOWLOG               SCAN 전환 + 3노드 격리            GET worst 15.5ms 해소",sz=8,c=WHITE,sp=1)
pa(tf,"  Replication           -                           Primary 50ops/s, Replica 200ops/s  Lag 0~1초, GTID 비동기",sz=8,c=WHITE,sp=1)
pa(tf,"",sz=3)
pa(tf,"검색 품질",sz=9,c=ACCENT,b=True,sp=3)
pa(tf,"  P@10                  0.827                       0.853 (+3.2%p)                     PhraseQuery(slop=2)+FeatureField",sz=8,c=WHITE,sp=2)
pa(tf,"  NDCG@10               0.6910                      0.7387 (+4.8%p, 5-Fold CV)         XGBoost LambdaMART 14피처",sz=8,c=WHITE,sp=1)
pa(tf,"  AI 검색               영문 Ai 문서만              '인공지능' 1위                      DB 동의어 쿼리 타임 확장",sz=8,c=WHITE,sp=1)
pa(tf,"  오타 교정             결과 0건                    '프로그래밍' 제안                   DirectSpellChecker 편집거리 2",sz=8,c=WHITE,sp=1)
pa(tf,"  Snippet               앞 150자 (무관)             검색어 주변 맥락                    UnifiedHighlighter 500자 StoredField",sz=8,c=WHITE,sp=1)
pa(tf,"  Facet                 없음                        30개 카테고리 전체 집계             SortedSetDocValuesFacetCounts",sz=8,c=WHITE,sp=1)
pa(tf,"  카테고리 분류         namespace (97% 편중)         28개 주제별 (83% 정확도)            키워드 기반 배치 분류, 90건 검증",sz=8,c=WHITE,sp=1)
pa(tf,"  금칙어                없음                        16,090개 (KO 3,094+EN 12,996)      Aho-Corasick O(N+Z)",sz=8,c=WHITE,sp=1)
pa(tf,"  AI 요약               없음                        RAG + SSE + 출처 링크              Gemini 15 RPM, 캐시 30분 TTL",sz=8,c=WHITE,sp=1)
pa(tf,"  검색 최종 (Facet+태그) -                           548ms (P95 2.61s)                  인덱스 42GB, Nori 사전 158K",sz=8,c=WHITE,sp=1)
pa(tf,"  LTR ON 한계           -                           3,088ms (72배 악화)                2코어 ARM CPU 포화 -> OFF",sz=8,c=RED,sp=1)
pa(tf,"",sz=3)
pa(tf,"인프라 규모: Lucene 42GB(5seg) | 재색인 ~2h | Nori 158,539 | Redis 3노드(4,620키) | HikariCP 5+15 | Kafka KRaft+DLQ",sz=7,c=LG,sp=3)
pa(tf,"기술 블로그: 26편 시리즈 (검색 인프라 10편 + 분산 아키텍처 8편 + 검색 품질 5편 + 참고 3편)",sz=7,c=LG,sp=2)

prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
